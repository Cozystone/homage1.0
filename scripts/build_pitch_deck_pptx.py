# -*- coding: utf-8 -*-
"""ATANOR pitch deck v2 (.pptx) — · .

 ( + ):
 · #050506 · (Helvetica World ) · #d2521f→#ff8a00
 · (constellation) — + = ' ' 
 · (XML), (XML gradFill), (#0d0e12++ )
 · · · 
Canva . 14, .
"""
import os, random, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

ASSETS=r"C:\0.ASKIM ALL-VIN\ATANOR-live-selfhood-scheduler"
OUTDIR=r"C:\0.ASKIM ALL-VIN\27., ATANOR DEMO"
LOGO=os.path.join(ASSETS,"assets","atanor_logo.png")

BG=RGBColor(0x05,0x05,0x06); CARD=RGBColor(0x0D,0x0E,0x12)
WHITE=RGBColor(0xFF,0xFF,0xFF); BODY=RGBColor(0x9A,0xA1,0xB2); DIM=RGBColor(0x6B,0x72,0x80)
ORANGE=RGBColor(0xD2,0x52,0x1F); ORANGE2=RGBColor(0xFF,0x8A,0x00); AMBER=RGBColor(0xFF,0x9F,0x1C)
HAIR=RGBColor(0x1D,0x23,0x30); LINEC=RGBColor(0x2A,0x30,0x3E)

FONT="Helvetica World"
W,H=Inches(13.333),Inches(7.5)
TOTAL=14

prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]
rng=random.Random(11)
_page=[0]

# ── low-level helpers ──────────────────────────────────────────────────────────
def _set_alpha(color_elem,pct):
    """srgbClr alpha(0~100%) ."""
    a=color_elem.makeelement(qn('a:alpha'),{'val':str(int(pct*1000))})
    color_elem.append(a)

def shape_alpha(sh,pct):
    """ (pct= %)."""
    sf=sh.fill._xPr.find(qn('a:solidFill'))
    if sf is None: return
    srgb=sf.find(qn('a:srgbClr'))
    if srgb is not None: _set_alpha(srgb,pct)

def grad_fill(sh,c1,c2,angle=0):
    """ (c1→c2)."""
    spPr=sh.fill._xPr
    for tag in ('a:solidFill','a:noFill','a:gradFill','a:blipFill','a:pattFill'):
        for e in spPr.findall(qn(tag)): spPr.remove(e)
    g=spPr.makeelement(qn('a:gradFill'),{})
    lst=g.makeelement(qn('a:gsLst'),{})
    for pos,c in ((0,c1),(100000,c2)):
        gs=g.makeelement(qn('a:gs'),{'pos':str(pos)})
        clr=g.makeelement(qn('a:srgbClr'),{'val':'%02X%02X%02X'%(c[0],c[1],c[2])})
        gs.append(clr); lst.append(gs)
    g.append(lst)
    lin=g.makeelement(qn('a:lin'),{'ang':str(angle*60000),'scaled':'1'})
    g.append(lin)
    ln=spPr.find(qn('a:ln'))
    spPr.insert(list(spPr).index(ln) if ln is not None else len(list(spPr)),g)

def slide():
    s=prs.slides.add_slide(BLANK); _page[0]+=1
    s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    return s

def _set_font(run,size,color,bold=True,name=FONT):
    f=run.font; f.name=name; f.size=Pt(size); f.bold=bold; f.color.rgb=color
    rPr=run._r.get_or_add_rPr()
    for tag in ("a:ea","a:cs"):
        for e in rPr.findall(qn(tag)): rPr.remove(e)
        rPr.append(rPr.makeelement(qn(tag),{"typeface":name}))

def text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_spacing=1.0,space_after=0):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    paras=runs if isinstance(runs[0],list) else [runs]
    for i,pruns in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=line_spacing
        if space_after: p.space_after=Pt(space_after)
        for t,size,color,bold in pruns:
            r=p.add_run(); r.text=t; _set_font(r,size,color,bold)
    return tb

def rect(s,x,y,w,h,color,line=None,alpha=None,round_=False):
    shp=MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sh=s.shapes.add_shape(shp,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=color
    if alpha is not None: shape_alpha(sh,alpha)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(0.75)
    sh.shadow.inherit=False
    return sh

def hairline(s,x,y,w,color=HAIR,pt=1.0):
    return rect(s,x,y,w,Pt(pt),color)

def grad_bar(s,x,y,w,h=Pt(3)):
    sh=rect(s,x,y,w,h,ORANGE); grad_fill(sh,ORANGE,AMBER,0); return sh

def dot(s,x,y,d,color,alpha=None):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,x,y,d,d)
    sh.fill.solid(); sh.fill.fore_color.rgb=color
    if alpha is not None: shape_alpha(sh,alpha)
    sh.line.fill.background(); sh.shadow.inherit=False
    return sh

def line_seg(s,x1,y1,x2,y2,color=LINEC,pt=0.75,alpha=None):

    x1,y1,x2,y2=(Emu(int(v)) for v in (x1,y1,x2,y2))
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2)
    cn.line.color.rgb=color; cn.line.width=Pt(pt)
    if alpha is not None:
        ln=cn.line._get_or_add_ln()
        sf=ln.find(qn('a:solidFill'))
        if sf is not None:
            srgb=sf.find(qn('a:srgbClr'))
            if srgb is not None: _set_alpha(srgb,alpha)
    cn.shadow.inherit=False
    return cn

def constellation(s,cx,cy,spread_x,spread_y,n=18,link=0.55,seedshift=0):
    """ — ."""
    r2=random.Random(23+seedshift)
    pts=[]
    for _ in range(n):
        x=cx+Emu(int(r2.gauss(0,spread_x))); y=cy+Emu(int(r2.gauss(0,spread_y)))
        if Emu(0)<x<W-Inches(0.15) and Emu(0)<y<H-Inches(0.15):
            pts.append((x,y))

    for i,(x1,y1) in enumerate(pts):
        dists=sorted(((math.hypot(int(x1-x2),int(y1-y2)),j) for j,(x2,y2) in enumerate(pts) if j!=i))
        for dctr,(dd,j) in enumerate(dists[:2]):
            if j>i and dd<int(spread_x)*1.4 and r2.random()<link:
                x2,y2=pts[j]
                line_seg(s,x1,y1,x2,y2,LINEC,0.6,alpha=55)
    for (x,y) in pts:
        dsz=Pt(r2.choice([2.2,2.6,3.2,4,5,6.5]))
        c=r2.choice([ORANGE,ORANGE2,AMBER,LINEC,RGBColor(0x8A,0x50,0x28)])
        dot(s,x,y,dsz,c,alpha=r2.choice([100,85,70,55]))

def progress(s):
    """ ."""
    n=TOTAL; gap=Inches(0.16); total_w=gap*(n-1)
    x0=W/2-total_w/2; y=H-Inches(0.34)
    for i in range(n):
        cur=(i==_page[0]-1)
        d=Pt(4.5 if cur else 3)
        dot(s,x0+gap*i-(d/2),y,d,ORANGE2 if cur else HAIR,alpha=100 if cur else 80)

def chrome(s,kicker):
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO,Inches(0.55),Inches(0.42),height=Inches(0.30))
    text(s,W-Inches(4.05),Inches(0.42),Inches(3.5),Inches(0.35),
         [(kicker,12,DIM,True)],align=PP_ALIGN.RIGHT)
    hairline(s,Inches(0.55),H-Inches(0.62),W-Inches(1.1))
    text(s,Inches(0.55),H-Inches(0.55),Inches(8),Inches(0.3),
         [("ATANOR · 2026.07 — 수치는 자체 실측, 커밋 단위 재현",9,DIM,False)])
    progress(s)

def head(s,kicker,big,sub=None,big_size=40):
    grad_bar(s,Inches(0.55),Inches(1.14),Inches(0.62),Pt(3.5))
    text(s,Inches(0.55),Inches(1.26),Inches(3.6),Inches(0.35),[(kicker,13,ORANGE2,True)])
    text(s,Inches(0.55),Inches(1.60),W-Inches(1.1),Inches(1.1),[(big,big_size,WHITE,True)],line_spacing=1.02)
    if sub:
        text(s,Inches(0.57),Inches(2.50),W-Inches(1.2),Inches(0.5),[(sub,15,BODY,False)])

def stat_cells(s,cells,y,cols=3,row_h=Inches(1.95)):
    gw=(W-Inches(1.1))/cols
    for i,(num,label,note) in enumerate(cells):
        x=Inches(0.55)+gw*(i%cols); yy=y+row_h*(i//cols)
        grad_bar(s,x,yy,Inches(0.5))
        text(s,x,yy+Inches(0.10),gw-Inches(0.35),Inches(1.0),[(num,46,WHITE,True)])
        text(s,x,yy+Inches(1.00),gw-Inches(0.35),Inches(0.4),[(label,14,BODY,True)])
        text(s,x,yy+Inches(1.35),gw-Inches(0.35),Inches(0.55),[(note,11,DIM,False)],line_spacing=1.05)

def bullets(s,items,y,size=15,gap=0.52,x=Inches(0.55),w=None):
    w=w or (W-Inches(1.1))
    for i,t in enumerate(items):
        yy=y+Inches(gap)*i
        dot(s,x,yy+Inches(0.09),Pt(5),ORANGE2)
        text(s,x+Inches(0.28),yy,w-Inches(0.28),Inches(0.6),[(t,size,BODY,False)])

def card(s,x,y,w,h,accent=ORANGE):
    c=rect(s,x,y,w,h,CARD,line=HAIR,round_=False)
    bar=rect(s,x,y,w,Pt(3),accent); grad_fill(bar,accent,AMBER,0)
    return c


s=slide()
constellation(s,Inches(10.4),Inches(2.3),Inches(1.9),Inches(1.5),n=26,seedshift=1)
constellation(s,Inches(2.2),Inches(6.1),Inches(1.2),Inches(0.8),n=12,seedshift=2)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO,W/2-Inches(1.7),Inches(1.95),width=Inches(3.4))
text(s,Inches(1.0),Inches(3.9),W-Inches(2.0),Inches(0.5),
     [("그래프 네이티브 로컬 AI — 프로그램 · 브라우저 · OS",17,BODY,False)],align=PP_ALIGN.CENTER)
gb=grad_bar(s,W/2-Inches(0.95),Inches(4.5),Inches(1.9),Pt(3))
text(s,Inches(1.0),Inches(4.72),W-Inches(2.0),Inches(0.65),
     [("환각하지 않고, GPU 없이, 출처를 증명하며 답하는 AI",22,WHITE,True)],align=PP_ALIGN.CENTER)
text(s,Inches(1.0),Inches(6.42),W-Inches(2.0),Inches(0.4),
     [("피치덱 · 2026.07 — 김안석 · 넥스트챌린지스쿨",12,DIM,False)],align=PP_ALIGN.CENTER)
progress(s)


s=slide(); chrome(s,"00 — 시작한 질문")
constellation(s,Inches(11.9),Inches(5.6),Inches(1.0),Inches(0.9),n=12,seedshift=3)
text(s,Inches(0.55),Inches(1.30),Inches(4.0),Inches(0.35),[("이 프로젝트가 시작된 곳",13,ORANGE2,True)])
grad_bar(s,Inches(0.55),Inches(1.18),Inches(0.62),Pt(3.5))
q1=card(s,Inches(0.55),Inches(1.85),W-Inches(1.65),Inches(1.55))
text(s,Inches(0.85),Inches(2.05),W-Inches(2.2),Inches(0.7),
     [("“지금의 AI는 왜 이렇게 비효율적인가?”",27,WHITE,True)])
text(s,Inches(0.85),Inches(2.78),W-Inches(2.2),Inches(0.5),
     [("질문 하나에 데이터센터 GPU가 수백 W를 태우고, 같은 질문에 매번 처음부터 다시 생각한다.",13.5,BODY,False)])
q2=card(s,Inches(0.55),Inches(3.65),W-Inches(1.65),Inches(1.55),accent=ORANGE2)
text(s,Inches(0.85),Inches(3.85),W-Inches(2.2),Inches(0.7),
     [("“ChatGPT-5.5급 AI를, 내 PC에서 직접 만들 수는 없을까?”",27,WHITE,True)])
text(s,Inches(0.85),Inches(4.58),W-Inches(2.2),Inches(0.5),
     [("모델을 빌리는 대신 — 지식·기억·추론·표현을 분리한 구조를 밑바닥부터.",13.5,BODY,False)])
text(s,Inches(0.55),Inches(5.75),W-Inches(1.1),Inches(0.55),
     [("그 두 질문의 답이, ",20,BODY,False),("ATANOR",20,ORANGE2,True),("다.",20,BODY,False)])


s=slide(); chrome(s,"01 — 문제")
head(s,"문제","우리는 AI를 소유하지 못하고, 빌려 쓴다",
     "중앙 클라우드 LLM 구조에는 갚을 수 없는 네 가지 청구서가 따라온다")
cells=[("환각","출처를 증명 못 하는 답","의료 · 법률 · 금융에 치명적"),
       ("유출","내 기억이 남의 서버로","데이터 주권의 상실"),
       ("GPU","질문마다 수백 W","개인이 평생 못 돌리는 비용"),
       ("망각","대화가 끝나면 잊는다","자라나는 AI 구조의 부재")]
gw=(W-Inches(1.1)-Inches(0.45))/4
for i,(num,label,note) in enumerate(cells):
    x=Inches(0.55)+(gw+Inches(0.15))*i
    card(s,x,Inches(3.05),gw,Inches(2.35))
    text(s,x+Inches(0.2),Inches(3.3),gw-Inches(0.4),Inches(0.8),[(num,34,WHITE,True)])
    text(s,x+Inches(0.2),Inches(4.15),gw-Inches(0.4),Inches(0.6),[(label,13,BODY,True)],line_spacing=1.1)
    text(s,x+Inches(0.2),Inches(4.75),gw-Inches(0.4),Inches(0.6),[(note,10.5,DIM,False)],line_spacing=1.1)
text(s,Inches(0.55),Inches(5.85),W-Inches(1.1),Inches(0.6),
     [("이건 ‘더 큰 모델’이 아니라, ",18,BODY,False),("‘다른 구조’",18,ORANGE2,True),("로 풀리는 문제다.",18,BODY,False)])


s=slide(); chrome(s,"02 — 해법")
constellation(s,Inches(12.1),Inches(5.3),Inches(0.9),Inches(0.9),n=10,seedshift=4)
head(s,"해법","뼈와 살 — 그리고 정직",
     "지식(뼈)은 그래프에, 말(살)은 검증된 문장에서 — 사실과 표현이 분리되어 오염되지 않는다")
bullets(s,[
 "답한다, 그러나 지어내지 않는다 — 근거가 없으면 ‘모른다’가 정답",
 "모든 답에 추론 증명서(근거 개념 · 도출 경로 · 보증) 부착 → 감사 가능",
 "후보 ↔ 검증 분리 — 새 지식은 게이트(출처·중복·모순·품질)를 통과해야 승격",
 "아이가 다섯 살까지 듣는 말로 언어를 깨치듯 — 순도 높은 15만 문장 식단",
 "개인 기억은 기기 밖으로 절대 나가지 않는다 (로컬-first)",
],Inches(3.15),gap=0.58)
q=card(s,Inches(0.55),Inches(6.05),Inches(7.2),Inches(0.72),accent=ORANGE2)
text(s,Inches(0.85),Inches(6.2),Inches(6.8),Inches(0.5),[("“진실이 커버리지보다 먼저다.”",19,WHITE,True)])


s=slide(); chrome(s,"03 — 제품")
head(s,"제품 · 2026.07","스스로 배우고, 스스로 지키며 돈다",
     "구상이 아니라 — 사람이 자리를 비워도 스스로를 지키며 달리는 시스템")
stat_cells(s,[("24h","자율 학습","위키·웹 검증 파이프 · 시간당 ~1,500문장"),
              ("15분","품질 파수꾼 주기","회귀 감지 → 학습 자동 동결 → 자동 복구"),
              ("8/10","프런티어 추론","은유 해석 · 반사실 · 유추 · 모순 감지")],Inches(3.15))
bullets(s,[
 "갭 자동 시딩 — 모르는 개념을 만나면 스스로 학습 대기열에 올린다",
 "연속 자아 · 호르몬 동역학 · 밤의 의회 — 파티클 아바타 ‘아토’ · 실시간 OPS 대시보드",
],Inches(5.35),size=13,gap=0.44)
text(s,Inches(0.55),Inches(6.35),Inches(11.5),Inches(0.5),
     [("목표는 챗봇이 아니라 — ",16,BODY,False),("프로그램 안에 사는 정직한 생명체.",16,ORANGE2,True)])


s=slide(); chrome(s,"04 — 측정된 사실")
head(s,"측정된 사실","마케팅이 아니라, 코드로 재현되는 숫자",
     "65문항 완성 게이트 배터리 · 봉인 홀드아웃 · 92,639 LOC · 커밋 979")
stat_cells(s,[("완성","배터리 판정 (65문항)","P0 23/23 · P1 31/32 · p50 ~2초"),
              ("92%","봉인 홀드아웃 QA","날조(근거 없는 단정) 0건"),
              ("94%","정직성 배터리","모르는 걸 아는 척하지 않는가"),
              ("2,600만","지식 그래프 트리플","적재 96만 행/s · 터보 3.0M 행/s"),
              ("0 GPU","답변 시 모델 추론","~0.001 Wh/질문 · 설치 ~11 MB"),
              ("0.25초","전 그래프 VRAM 미러","606MB · 64차원 임베딩 자체 훈련")],Inches(2.95),row_h=Inches(1.88))
text(s,Inches(0.55),Inches(6.50),W-Inches(1.1),Inches(0.35),
     [("정직한 한계: 커버리지는 LLM보다 좁고 장문 창작은 아직 얕다(식단 축적 중) — 숨긴 결함이 아니라 설계상 트레이드오프.",10,DIM,False)])


s=slide(); chrome(s,"05 — 3박자")
head(s,"장기 비전 Ⅰ","하나의 그래프 뇌가, 세 겹의 몸을 입는다",
     "LLM은 클라우드에서 아래로 — ATANOR는 기기에서 위로")
gw=(W-Inches(1.1)-Inches(0.5))/3
cards=[("① 프로그램","지금 · 완성 게이트 통과","로컬 엔진+웹앱. 65문항 배터리 ‘완성’ 판정 — 실사용 수준 검증",ORANGE),
       ("② 브라우저","진행 · 확장 v0.8.3","AI가 스스로 서핑하며 페이지를 증류해 배우는 ‘인지 서핑’ → AI 브라우저",ORANGE2),
       ("③ OS","첫 부팅 완료","Rust 자체 컴포지터(M1) — ATANOR Linux 부팅 실증 → 파티클 셸 OS",AMBER)]
cy=Inches(3.1)
for i,(t,st,d,ac) in enumerate(cards):
    x=Inches(0.55)+(gw+Inches(0.25))*i
    card(s,x,cy,gw,Inches(2.35),accent=ac)
    text(s,x+Inches(0.25),cy+Inches(0.22),gw-Inches(0.5),Inches(0.55),[(t,21,WHITE,True)])
    text(s,x+Inches(0.25),cy+Inches(0.80),gw-Inches(0.5),Inches(0.4),[(st,12.5,ORANGE2,True)])
    text(s,x+Inches(0.25),cy+Inches(1.20),gw-Inches(0.5),Inches(1.0),[(d,11.5,BODY,False)],line_spacing=1.15)

mid=cy+Inches(1.15)
line_seg(s,Inches(0.55)+gw,mid,Inches(0.55)+gw+Inches(0.25),mid,ORANGE2,1.4)
line_seg(s,Inches(0.55)+gw*2+Inches(0.25),mid,Inches(0.55)+gw*2+Inches(0.5),mid,ORANGE2,1.4)
bullets(s,[
 "각 층이 같은 뇌를 공유 — 배울수록 셋이 함께 똑똑해진다",
 "개인 데이터가 앱→브라우저→OS 어디서도 밖으로 새지 않는 유일한 스택 — 구조적 해자",
],Inches(5.75),size=13.5,gap=0.46)


s=slide(); chrome(s,"05+ — 피지컬 AI")
constellation(s,Inches(11.3),Inches(2.2),Inches(1.5),Inches(1.2),n=22,seedshift=5)
head(s,"장기 비전 Ⅱ","잊지 않는 뇌는, 피지컬 AI의 핵심이 된다",
     "맥락을 사람처럼 이해하고 잊지 않는 로컬 뇌 — 3박자 다음의 네 번째 몸",big_size=34)
bullets(s,[
 "스마트 글래스 · 웨어러블 · 로봇에 연결 — 개인의 모든 맥락(본 것 · 들은 것 · 한 일)을 기기 안에서 기억",
 "클라우드 왕복 없음 — 실시간성과 프라이버시를 동시에 (피지컬 AI의 필요조건)",
 "기기가 바뀌어도 — 같은 뇌, 같은 자아, 같은 기억 (연속 자아 실증 완료)",
 "시장이 마중 나온다 — AI 글래스 판매 2026년 한 해 4배(600만→2,000만 대 전망) · 피지컬 AI $81.6B→$960B(2033 추정)",
],Inches(3.15),size=14,gap=0.6)
q=card(s,Inches(0.55),Inches(5.95),Inches(10.6),Inches(0.72),accent=ORANGE2)
text(s,Inches(0.85),Inches(6.1),Inches(10.1),Inches(0.5),
     [("클라우드 LLM은 물리 세계의 매 순간을 감당할 수 없다 — ",15.5,BODY,False),("로컬 뇌만이 할 수 있다.",15.5,ORANGE2,True)])


s=slide(); chrome(s,"06 — 시장")
head(s,"시장","틀리면 안 되고, 나가면 안 되고, 싸야 하는 곳",
     "‘가장 똑똑한 모델’ 경쟁이 아니라 — 신뢰 세그먼트를 정조준",big_size=36)
stat_cells(s,[("$11.8B","엣지/온디바이스 AI · 2025","→ $56.8B (2030) · CAGR 36.9% (BCC)"),
              ("$960B","피지컬 AI · 2033 (추정)","$81.6B(2025)에서 CAGR 36.1% (Grand View)"),
              ("규제","EU AI Act — 감사가능성","‘출처를 증명하는 AI’ = 규제 대응 자산")],Inches(3.15))
bullets(s,[
 "초기: 개발자·연구자 (자기 문서·코드를 장기 기억으로) → 소규모 팀 온프레미스",
 "확장: 금융 · 의료 · 법률 · 공공 — 환각과 유출이 허용되지 않는 산업",
],Inches(5.55),size=14,gap=0.48)


s=slide(); chrome(s,"07 — 비즈니스 모델")
head(s,"비즈니스 모델","COGS가 0에 수렴하는 AI",
     "추론에 GPU가 들지 않는다 — 가격 · 마진 · 엣지 배포 전부의 구조적 우위")
rows=[("개인 구독","월 9,900–29,000원","로컬 그래프 관리 · 백그라운드 학습"),
      ("Pro (개발자·연구자)","월 39,000–99,000원","문서/코드 인덱싱 · 프로젝트 브레인"),
      ("팀/기업 온프레미스","구축 + 월유지 + 시트","데이터 반출 없는 사내 지식 AI"),
      ("Graph Cartridge 마켓","수수료 20–30%","도메인 그래프(법률·의료·창업) 장터"),
      ("브라우저 → OS → 피지컬","무료 퍼널 → 라이선스","3박자 완성 시 기기·OS·피지컬 뇌 (장기)")]
y=Inches(3.05)
for i,(a,bp,d) in enumerate(rows):
    yy=y+Inches(0.66)*i
    hairline(s,Inches(0.55),yy,W-Inches(1.1))
    dot(s,Inches(0.55),yy+Inches(0.20),Pt(5),ORANGE2)
    text(s,Inches(0.82),yy+Inches(0.08),Inches(3.3),Inches(0.5),[(a,14.5,WHITE,True)])
    text(s,Inches(4.35),yy+Inches(0.09),Inches(3.2),Inches(0.5),[(bp,13.5,ORANGE2,True)])
    text(s,Inches(7.7),yy+Inches(0.11),W-Inches(8.25),Inches(0.5),[(d,12,BODY,False)])
text(s,Inches(0.55),Inches(6.55),Inches(11.5),Inches(0.4),
     [("검증 순서: Pro → 온프레미스 → 카트리지 마켓 → 3박자 스택",13.5,DIM,False)])


s=slide(); chrome(s,"08 — 경쟁")
head(s,"경쟁","방향이 반대다",
     "그들은 클라우드 모델을 아래로 내린다 — 우리는 로컬 그래프 뇌를 위로 올린다")
rows=[("ChatGPT · Claude","범용 추론의 정점","로컬 지식 소유 없음 · 검증 그래프 아님"),
      ("Perplexity","검색 + 출처 답변","검색 응답일 뿐 — 지식이 쌓이지 않는다"),
      ("로컬 LLM (Llama 등)","오프라인 구동","환각 여전 · 수 GB 가중치 — ATANOR는 0"),
      ("AI 브라우저 (Arc·Comet)","브라우징 보조","클라우드 LLM을 얹음 — 우리는 로컬 뇌가 브라우저를 쓴다")]
y=Inches(3.1)
for i,(a,st,d) in enumerate(rows):
    yy=y+Inches(0.72)*i
    hairline(s,Inches(0.55),yy,W-Inches(1.1))
    text(s,Inches(0.55),yy+Inches(0.09),Inches(3.5),Inches(0.5),[(a,14.5,WHITE,True)])
    text(s,Inches(4.25),yy+Inches(0.10),Inches(2.9),Inches(0.5),[(st,12.5,BODY,False)])
    text(s,Inches(7.25),yy+Inches(0.10),W-Inches(7.8),Inches(0.6),[(d,12,DIM,False)])
q=card(s,Inches(0.55),Inches(6.0),Inches(11.4),Inches(0.72),accent=ORANGE)
text(s,Inches(0.85),Inches(6.15),Inches(11.0),Inches(0.5),
     [("해자 = ",15.5,BODY,False),("정직성 + 경량성 + 프라이버시 + 3박자 수직통합",15.5,ORANGE2,True),("의 결합",15.5,BODY,False)])


s=slide(); chrome(s,"09 — 로드맵")
head(s,"로드맵","식단 완주 → 베타 → 온프레미스 → OS → 피지컬",big_size=32)
cols=[("단기 · 0–6개월",["문장 식단 15만 완주 → 유창성 개화","클로즈드 베타 100명 + 사전예약","공개 벤치마크 · 에너지 백서"],ORANGE),
      ("중기 · 6–18개월",["온프레미스 패키지 + 카트리지","브라우저 확장 정식 배포 (스토어)","특허 출원 · Atlas Network v1"],ORANGE2),
      ("장기 · 18개월+",["ATANOR OS — 디바이스 연속성","피지컬 AI 뇌 (글래스·웨어러블)","'믿을 수 있는 답'의 표준"],AMBER)]
gw=(W-Inches(1.1)-Inches(0.5))/3

line_seg(s,Inches(0.7),Inches(3.02),W-Inches(0.7),Inches(3.02),LINEC,1.2)
for i,(t,items,ac) in enumerate(cols):
    x=Inches(0.55)+(gw+Inches(0.25))*i
    dot(s,x+Inches(0.06),Inches(2.93),Pt(9),ac)
    text(s,x,Inches(3.2),gw,Inches(0.45),[(t,16.5,WHITE,True)])
    hairline(s,x,Inches(3.66),gw-Inches(0.15))
    yy=Inches(3.85)
    for it in items:
        dot(s,x,yy+Inches(0.08),Pt(4.5),ac)
        text(s,x+Inches(0.22),yy,gw-Inches(0.3),Inches(0.75),[(it,12,BODY,False)],line_spacing=1.12)
        yy+=Inches(0.78)


s=slide(); chrome(s,"10 — 팀")
constellation(s,Inches(12.0),Inches(5.5),Inches(0.9),Inches(0.8),n=10,seedshift=6)
head(s,"팀","김안석 — 단독 창업자",
     "넥스트챌린지스쿨 (서울시교육청 인가 · 국내 첫 창업 특화 대안고) — Google · Intel · Thales 협력")
bullets(s,[
 "시나브로 (문학인 SNS) — 학생창업유망팀 U300+ 도약트랙 최종선발",
 "WETHUS (학생 창업 플랫폼) — ‘모두의 창업’ 1차 심사 통과",
 "BELIFE (개인 인텔리전스) — 생각 · 감정 · 가치관 구조화 실험",
 "본 프로젝트 실측 — 6개월 · 커밋 979 · 92,639 LOC · 78패키지 · 테스트 606 (그래프 스토어→형태소 엔진→Rust OS까지 단독)",
],Inches(3.1),size=13.5,gap=0.5)
text(s,Inches(0.55),Inches(5.35),W-Inches(1.1),Inches(0.6),
     [("“이 사람들의 생각과 기억은, 결국 누구의 것인가.”",22,WHITE,True)])
text(s,Inches(0.55),Inches(6.15),W-Inches(1.1),Inches(0.4),
     [("고등학생이지만 — 아이디어가 아니라 측정 가능한 사실까지 만들어내는 실행력.",13,DIM,False)])


s=slide()
constellation(s,Inches(2.1),Inches(5.7),Inches(1.4),Inches(1.0),n=16,seedshift=7)
constellation(s,Inches(11.3),Inches(1.7),Inches(1.2),Inches(0.9),n=14,seedshift=8)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO,W/2-Inches(1.0),Inches(1.8),width=Inches(2.0))
text(s,Inches(1.0),Inches(3.3),W-Inches(2.0),Inches(0.85),
     [("답한다, 그러나 지어내지 않는다.",36,WHITE,True)],align=PP_ALIGN.CENTER)
text(s,Inches(1.0),Inches(4.3),W-Inches(2.0),Inches(0.5),
     [("AI를 빌리는 시대에서, 소유하는 시대로 — 프로그램 · 브라우저 · OS, 그리고 피지컬.",15,BODY,False)],align=PP_ALIGN.CENTER)
grad_bar(s,W/2-Inches(0.95),Inches(5.0),Inches(1.9),Pt(3))
text(s,Inches(1.0),Inches(5.25),W-Inches(2.0),Inches(0.4),
     [("김안석 · github.com/Cozystone/ATANOR · 2026.07",12,DIM,False)],align=PP_ALIGN.CENTER)
text(s,Inches(1.0),Inches(6.6),W-Inches(2.0),Inches(0.35),
     [("본 덱의 모든 수치는 자체 실측이며 커밋 단위로 재현됩니다.",10,DIM,False)],align=PP_ALIGN.CENTER)
progress(s)

out=os.path.join(OUTDIR,"ATANOR_피치덱_v2.pptx")
prs.save(out)
print("WROTE",out,"| slides:",sum(1 for _ in prs.slides),"| logo:",os.path.exists(LOGO))
